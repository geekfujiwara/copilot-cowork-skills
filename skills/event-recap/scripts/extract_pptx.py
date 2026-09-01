#!/usr/bin/env python3
"""
extract_pptx.py — 当日資料 PPTX から全テキスト（本文・表・ノート）を抽出する。

汎用イベント開催レポート用。どの案件でも当日資料スライドのファクト抽出に使える。

使い方:
    python extract_pptx.py <pptx_path> [--out <text_path>]

出力: スライドごとに区切ったプレーンテキスト。--out 未指定なら標準出力。
依存: python-pptx（コンテナにプリインストール済み。pip install しないこと）
"""
import sys
import glob
import argparse


def extract(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"===== SLIDE {i} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    out.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                out.append("[NOTES] " + n)
    return "\n".join(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", help="PPTX ファイルパス。glob パターン可（最初の一致を使用）")
    ap.add_argument("--out", help="出力先テキストファイル")
    args = ap.parse_args()

    matches = glob.glob(args.pptx)
    path = matches[0] if matches else args.pptx
    text = extract(path)

    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"FILE: {path}\nCHARS: {len(text)}\nWROTE: {args.out}")
    else:
        sys.stdout.write(text)


if __name__ == "__main__":
    main()
