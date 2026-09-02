#!/usr/bin/env python3
"""Create editable hierarchy/maturity diagrams with native PowerPoint shapes."""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

CW = 2_606_040
INNER_OFFSET = 182_880
IW = 2_240_280
BH = 274_320
BOXW = 685_800
GAP = 91_440
LINE_W = 9_525

COLORS = {
    "navy": "1E2761",
    "blue": "0F6CBD",
    "body": "3A4A6B",
    "muted": "616A82",
    "empty_fill": "F5F7FB",
    "empty_line": "C7DBF5",
    "growing_fill": "EAF1FB",
    "step1_line": "AEBCD4",
    "white": "FFFFFF",
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _text(shape, value: str, size: float, color: str, bold: bool = True) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Emu(30_000)
    frame.margin_top = frame.margin_bottom = Emu(12_000)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Yu Gothic UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _rect(slide, x: int, y: int, w: int, h: int, color: str):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    _fill(shape, color)
    _no_line(shape)
    return shape


def _rounded_box(slide, x: int, y: int, w: int, h: int, fill: str,
                 line: str, label: str = "", text_color: str = "1E2761",
                 dashed: bool = False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
    )
    if len(shape.adjustments):
        shape.adjustments[0] = 0.25
    _fill(shape, fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Emu(LINE_W)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if label:
        _text(shape, label, 8, text_color)
    return shape


def add_hierarchy_diagram(slide, card_x: int, y0: int, stage: int,
                          headquarters: str = "本部",
                          divisions: tuple[str, str, str] = ("事業部A", "事業部B", "事業部C"),
                          caption: str = "段階的に参画範囲を拡大") -> None:
    """Add one hierarchy diagram. Stage must be 1..4."""
    if stage not in {1, 2, 3, 4}:
        raise ValueError("stage must be 1, 2, 3, or 4")
    if len(divisions) != 3:
        raise ValueError("divisions must contain exactly three labels")

    inner = card_x + INNER_OFFSET
    parent_x = card_x + 525_780
    parent_w = 1_554_480
    parent_center = card_x + 1_303_020
    bar_y = y0 + BH + 91_440
    child_y = y0 + BH + 182_880
    child_centers = [inner + index * (BOXW + GAP) + BOXW // 2 for index in range(3)]
    connector = COLORS["step1_line"] if stage == 1 else COLORS["blue"]

    _rounded_box(
        slide, parent_x, y0, parent_w, BH,
        COLORS["navy"], COLORS["navy"], headquarters, COLORS["white"],
    )
    _rect(slide, parent_center - LINE_W // 2, y0 + BH, LINE_W, 91_440, connector)
    _rect(slide, child_centers[0], bar_y, child_centers[-1] - child_centers[0], LINE_W, connector)

    if stage == 1:
        states = ("empty", "empty", "empty")
    elif stage == 2:
        states = ("growing", "empty", "empty")
    elif stage == 3:
        states = ("active", "empty", "empty")
    else:
        states = ("active", "active", "active")

    for index, state in enumerate(states):
        line_color = COLORS["empty_line"] if state == "empty" else connector
        _rect(slide, child_centers[index] - LINE_W // 2, bar_y, LINE_W, 91_440, line_color)
        child_x = inner + index * (BOXW + GAP)
        if state == "empty":
            _rounded_box(
                slide, child_x, child_y, BOXW, BH,
                COLORS["empty_fill"], COLORS["empty_line"], dashed=True,
            )
        elif state == "growing":
            _rounded_box(
                slide, child_x, child_y, BOXW, BH,
                COLORS["growing_fill"], COLORS["empty_line"], divisions[index], COLORS["navy"],
            )
        else:
            _rounded_box(
                slide, child_x, child_y, BOXW, BH,
                COLORS["blue"], COLORS["blue"], divisions[index], COLORS["white"],
            )

    caption_box = slide.shapes.add_textbox(
        Emu(inner), Emu(y0 + BH * 2 + 228_600), Emu(IW), Emu(190_000)
    )
    _text(caption_box, caption, 8, COLORS["muted"], bold=False)


def build_sample(output: Path) -> None:
    """Create four stage cards and save once for a single render review."""
    presentation = Presentation()
    presentation.slide_width = Emu(12_192_000)
    presentation.slide_height = Emu(6_858_000)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb("F5F7FB")

    left = 457_200
    gap = 152_400
    card_y = 1_500_000
    card_h = 1_650_000
    for index in range(4):
        card_x = left + index * (CW + gap)
        card = _rounded_box(slide, card_x, card_y, CW, card_h, "FFFFFF", "E4E7EF")
        _text(card, "", 8, COLORS["body"], bold=False)
        label = slide.shapes.add_textbox(Emu(card_x), Emu(card_y + 110_000), Emu(CW), Emu(250_000))
        _text(label, f"STEP {index + 1}", 12, COLORS["navy"])
        add_hierarchy_diagram(
            slide, card_x, card_y + 480_000, index + 1,
            caption=("全て未参画", "1部門を育成", "1部門が稼働", "3部門へ展開")[index],
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", type=Path, default=Path("working/hierarchy-sample.pptx"))
    args = parser.parse_args()
    build_sample(args.out)
    print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
