#!/usr/bin/env python3
"""Create an editable long-tail analysis slide with native PowerPoint shapes."""
from __future__ import annotations

import argparse
import math
from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

SLIDE_W = 12_192_000
SLIDE_H = 6_858_000
MARGIN = 457_200
CARD_Y = 1_335_024
CARD_H = 3_008_376
CARD_W = SLIDE_W - 2 * MARGIN

COLORS = {
    "background": "F5F7FB",
    "card": "FFFFFF",
    "line": "E4E7EF",
    "navy": "1E2761",
    "blue": "0F6CBD",
    "muted": "616A82",
    "text": "3A4A6B",
    "light_blue": "8FB3E0",
    "pale_line": "C7DBF5",
    "pale_blue": "EAF1FB",
    "white": "FFFFFF",
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.fill.transparency = transparency


def _set_text(shape, text: str, size: float, color: str, *, bold: bool = False,
              align=PP_ALIGN.LEFT, rotation: int | None = None) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Emu(12_000)
    frame.margin_top = frame.margin_bottom = Emu(8_000)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if rotation is not None:
        frame.rotation = rotation
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Yu Gothic UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _textbox(slide, x: int, y: int, w: int, h: int, text: str, size: float,
             color: str, *, bold: bool = False, align=PP_ALIGN.LEFT,
             rotation: int | None = None):
    shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    _set_text(shape, text, size, color, bold=bold, align=align, rotation=rotation)
    return shape


def _rect(slide, x: int, y: int, w: int, h: int, fill: str,
          line: str | None = None, transparency: int = 0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    _fill(shape, fill, transparency)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _round_rect(slide, x: int, y: int, w: int, h: int, fill: str, line: str):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
    )
    if len(shape.adjustments):
        shape.adjustments[0] = 0.12
    _fill(shape, fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(0.75)
    return shape


def _line(slide, x1: int, y1: int, x2: int, y2: int, color: str,
          width: float = 1.0, dashed: bool = False):
    shape = slide.shapes.add_connector(1, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def _badge(slide, x: int, y: int, w: int, text: str) -> None:
    shape = _round_rect(slide, x, y, w, 260_000, COLORS["blue"], COLORS["blue"])
    _set_text(shape, text, 11, COLORS["white"], bold=True, align=PP_ALIGN.CENTER)


def _long_tail_curve(slide, x: int, y: int, w: int, h: int) -> None:
    """Draw one editable freeform polyline sampled from a long-tail curve."""
    points: list[tuple[int, int]] = []
    samples = 36
    for index in range(samples):
        t = index / (samples - 1)
        px = int(x + t * w)
        normalized = 0.08 + 0.92 * math.exp(-5.2 * t)
        py = int(y + h * (1 - normalized))
        points.append((px, py))
    builder = slide.shapes.build_freeform(Emu(points[0][0]), Emu(points[0][1]))
    builder.add_line_segments([(Emu(px), Emu(py)) for px, py in points[1:]], close=False)
    curve = builder.convert_to_shape()
    curve.fill.background()
    curve.line.color.rgb = _rgb(COLORS["blue"])
    curve.line.width = Pt(2.5)


def add_long_tail_diagram(
    slide,
    title_prefix: str,
    title_accent: str,
    subtitle: str,
    *,
    x_label: str = "対象（件数順）",
    y_label: str = "成果・影響度",
    left_badge: str = "ヘッド",
    right_badge: str = "ロングテール",
    left_note: str = "少数の重点対象を深く分析",
    right_note: str = "多数の小さな機会を継続的に捉える",
    left_card_title: str = "重点領域",
    left_card_body: str = "高い影響度の対象へ、個別施策と十分なリソースを配分する。",
    right_card_title: str = "拡張領域",
    right_card_body: str = "標準化と自動化により、小規模な機会を広く積み上げる。",
) -> None:
    """Add a complete long-tail slide to a 16:9 presentation."""
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(COLORS["background"])

    title = slide.shapes.add_textbox(Emu(MARGIN), Emu(228_600), Emu(CARD_W), Emu(502_920))
    frame = title.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    for value, color in ((title_prefix, COLORS["navy"]), (title_accent, COLORS["blue"])):
        run = paragraph.add_run()
        run.text = value
        run.font.name = "Yu Gothic UI"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = _rgb(color)

    _textbox(slide, MARGIN, 768_096, CARD_W, 380_000, subtitle, 14, COLORS["text"])
    _round_rect(slide, MARGIN, CARD_Y, CARD_W, CARD_H, COLORS["card"], COLORS["line"])

    plot_x = MARGIN + 730_000
    plot_y = CARD_Y + 360_000
    plot_w = CARD_W - 1_100_000
    plot_h = 1_650_000
    split_x = plot_x + int(plot_w * 0.38)

    _rect(slide, plot_x, plot_y, split_x - plot_x, plot_h, COLORS["pale_blue"])
    _rect(slide, split_x, plot_y, plot_x + plot_w - split_x, plot_h, COLORS["pale_line"])
    _line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, COLORS["muted"], 1.0)
    _line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, COLORS["muted"], 1.0)
    _line(slide, split_x, plot_y, split_x, plot_y + plot_h, COLORS["light_blue"], 1.0, dashed=True)
    _long_tail_curve(slide, plot_x + 30_000, plot_y + 60_000, plot_w - 60_000, plot_h - 100_000)

    _textbox(slide, plot_x + 30_000, plot_y + plot_h + 35_000, plot_w - 60_000, 220_000,
             x_label, 10, COLORS["muted"], align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN + 150_000, plot_y + 260_000, 300_000, 1_100_000,
             y_label, 10, COLORS["muted"], align=PP_ALIGN.CENTER, rotation=270)

    _badge(slide, plot_x + 150_000, plot_y + 110_000, 830_000, left_badge)
    _badge(slide, split_x + 220_000, plot_y + 590_000, 1_250_000, right_badge)
    _textbox(slide, plot_x + 120_000, plot_y + 420_000, split_x - plot_x - 240_000, 340_000,
             left_note, 10, COLORS["text"], bold=True)
    _textbox(slide, split_x + 210_000, plot_y + 920_000, plot_x + plot_w - split_x - 420_000, 340_000,
             right_note, 10, COLORS["text"], bold=True)

    footer_y = CARD_Y + CARD_H + 260_000
    footer_gap = 220_000
    footer_w = (CARD_W - footer_gap) // 2
    for index, (card_title, card_body) in enumerate((
        (left_card_title, left_card_body), (right_card_title, right_card_body)
    )):
        card_x = MARGIN + index * (footer_w + footer_gap)
        _round_rect(slide, card_x, footer_y, footer_w, 1_230_000, COLORS["card"], COLORS["line"])
        _textbox(slide, card_x + 210_000, footer_y + 150_000, footer_w - 420_000, 280_000,
                 card_title, 16, COLORS["navy"], bold=True)
        _textbox(slide, card_x + 210_000, footer_y + 500_000, footer_w - 420_000, 560_000,
                 card_body, 14, COLORS["text"])


def build_sample(output: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Emu(SLIDE_W)
    presentation.slide_height = Emu(SLIDE_H)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_long_tail_diagram(
        slide,
        "重点施策と",
        "ロングテール施策を両立する",
        "影響度の高い対象を深く支援しながら、標準化した施策で幅広い機会を積み上げます。",
    )
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", type=Path, default=Path("working/long-tail-sample.pptx"))
    args = parser.parse_args(argv)
    build_sample(args.out)
    print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
