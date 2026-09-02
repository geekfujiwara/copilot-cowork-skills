#!/usr/bin/env python3
"""Extract text and table content from a selected PPTX slide range."""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def extract(path: Path, start: int = 1, end: int | None = None) -> None:
    """Print non-empty text and table rows from the requested slides."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx が必要です") from exc

    presentation = Presentation(path)
    total = len(presentation.slides)
    last = end or total
    if last > total:
        last = total
    print(f"SLIDES: {total} (showing {start}..{last})")

    for index, slide in enumerate(presentation.slides, 1):
        if index < start or index > last:
            continue
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    line = " | ".join(cell.text.strip() for cell in row.cells)
                    if line.strip(" |"):
                        chunks.append("[TBL] " + line)
        if chunks:
            print(f"\n==== Slide {index} ====")
            print("\n".join(chunks))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="入力PPTXのパス")
    parser.add_argument("start", type=int, nargs="?", default=1, help="開始スライド")
    parser.add_argument("end", type=int, nargs="?", help="終了スライド")
    args = parser.parse_args(argv)
    if not args.pptx.is_file():
        parser.error(f"ファイルが見つかりません: {args.pptx}")
    if args.start < 1 or (args.end is not None and args.end < args.start):
        parser.error("スライド範囲が不正です")
    try:
        extract(args.pptx.resolve(), args.start, args.end)
    except (OSError, RuntimeError, ValueError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
