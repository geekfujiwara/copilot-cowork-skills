#!/usr/bin/env python3
"""
extract-deck-text.py — もと資料 PPTX の本文テキストを抽出する（汎用）。

使い方:
    python3 extract-deck-text.py "{{ABS_PATH_TO_PPTX}}" [START_SLIDE] [END_SLIDE]

ポイント（talk-prep 教訓）:
- ReadFileContent 等で取得した PPTX は grounding/downloads/ などに保存される。
- python-pptx には **絶対パス** を渡す（cwd 相対だと PackageNotFoundError になりうる）。
- text_frame と table の両方を拾う。巨大ファイルは START/END で分割して読む。
"""
import argparse
import sys
from pathlib import Path


def extract(path, start=1, end=None):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx が必要です") from exc
    prs = Presentation(path)  # 絶対パス推奨
    total = len(prs.slides)
    end = end or total
    print(f"SLIDES: {total} (showing {start}..{end})")
    for i, slide in enumerate(prs.slides, 1):
        if i < start or i > end:
            continue
        chunks = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if sh.has_table:
                for row in sh.table.rows:
                    line = " | ".join(c.text.strip() for c in row.cells)
                    if line.strip(" |"):
                        chunks.append("[TBL] " + line)
        if chunks:
            print(f"\n==== Slide {i} ====")
            print("\n".join(chunks))


def main(argv=None):
    parser = argparse.ArgumentParser(description="PPTX から本文と表のテキストを抽出します")
    parser.add_argument("pptx", type=Path, help="入力 PPTX のパス")
    parser.add_argument("start", type=int, nargs="?", default=1, help="開始スライド（1 始まり）")
    parser.add_argument("end", type=int, nargs="?", help="終了スライド（省略時は末尾）")
    args = parser.parse_args(argv)
    if not args.pptx.is_file():
        parser.error(f"ファイルが見つかりません: {args.pptx}")
    if args.start < 1 or (args.end is not None and args.end < args.start):
        parser.error("スライド範囲が不正です")
    try:
        extract(args.pptx, args.start, args.end)
    except (OSError, RuntimeError, ValueError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
